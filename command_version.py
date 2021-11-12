from requests import get
import zipfile
from os import path, rename, mkdir, chdir
from tqdm import tqdm
import re
import io
import asyncio
import aiohttp
import pptx

headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 6.1; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) '
                         'Chrome/90.0.4430.212 Safari/537.36'}
url_all = 'http://www.1ppt.com'

if not path.isdir('download'):  # 建立下载文件夹
    mkdir('download')
chdir('download')


def req_for_web(url):  # 该函数负责获取页面
    req1 = get(url, headers=headers)
    req1.encoding = 'gbk'
    return req1.text


def get_re(req_in, pattern):  # 该函数与上面的函数一起实现页面内容提取
    re1 = re.compile(pattern)
    result = re1.findall(req_in)
    return result


req_site = req_for_web(url_all)
all_choose_link = get_re(req_site, 'href=\"(.*?[^article|office])\" title=\".*?\"[^>]')
all_choose_title = get_re(req_site, 'href=\".*?[^article|office]\" title=\"(.*?)\"[^>]')
all_choose_link.append('/moban')
all_choose_title.append('全部')
title = [str(all_choose_title.index(x) + 1) + '、' + x for x in all_choose_title]
for i in title:
    print(i, end=' ')
print('')
choice = int(input('请选择:'))
if not path.isdir(all_choose_title[choice - 1]):
    mkdir(all_choose_title[choice - 1])
chdir(all_choose_title[choice - 1])  # 以上内容为分类功能的实现


def get_choice_link(choose_title):  # 合成分类链接
    choice_link1 = all_choose_link[choose_title - 1]
    return choice_link1


async def aio_get(url, pattern):  # 此函数实现异步提取PPT链接
    async with aiohttp.ClientSession() as r:
        async with r.get(url, headers=headers) as rep:
            i1 = await rep.text(encoding='gbk')
            re1 = re.compile(pattern)
            re2 = re1.findall(i1)
            return re2


def loop_aio(url_list, pattern):  # 异步实现
    task = [asyncio.ensure_future(aio_get(url1, pattern)) for url1 in url_list]
    tasks = asyncio.gather(*task)
    loop = asyncio.get_event_loop()
    loop.run_until_complete(tasks)
    return tasks.result()


def file_download(url):  # 获取真实链接
    web = req_for_web(url)
    file_first_link_list = get_re(web, 'a href=\"(.*?)\" .*</h2')
    link_2 = loop_aio([url_all + i1 for i1 in file_first_link_list], 'a href=\'(.*?)\' target=')
    link_d = loop_aio([url_all + l1[0] for l1 in link_2], 'a href=\"(.*?)\" target=')
    return link_d


def format_link(web, page):  # 处理翻页
    page_link = get_re(web, 'href=\"(.*?)\">2')[0]
    return '/' + page_link[:-6] + str(page) + '.html'


def get_the_max_page(web):  # 获取最大页码，避免for循环报错
    max_link = get_re(web, 'href=\'(.*?)\'>末页')[0]
    if '_' in max_link:
        return max_link.split('_')[-1][:-5]
    elif '/' in max_link:
        return max_link.split('/')[-1][:-5]


def replace_text(text_frame):  # 该函数实现的是ppt文本替换功能
    text_need_replace = [('第一PPT', ''), ('www.1ppt.com', ''), ('1ppt', '')]
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for tt in text_need_replace:
                if tt[0] in run.text:
                    run.text = run.text.replace(tt[0], tt[1])


def process_ppt(filename_open):  # 对下载完成的PPT进行处理
    prs = pptx.Presentation(filename_open)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # 判断Shape是否含有文本框
                text_frame = shape.text_frame
                replace_text(text_frame)  # 调用replace_text函数实现文本替换
            if shape.has_table:  # 判断Shape是否含有表格
                table = shape.table
                for cell in table.iter_cells():  # 遍历表格的cell
                    text_frame = cell.text_frame
                    replace_text(text_frame)  # 调用replace_text函数实现文本替换
    prs = pptx.Presentation(filename_open)
    slides = prs.slides
    rId = slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]
    prs.save(filename_open)


def auto_decode(name):  # 对文件名进行处理
    gbk_name = str(name).encode('cp437').decode('gbk')
    x = 1
    if '(' not in gbk_name:
        while True:
            if path.isfile(gbk_name):
                gbk_name = str(gbk_name).split('.')[0] + f'({x}).' + str(gbk_name).split('.')[1]
                x += 1
            else:
                break
    else:
        while True:
            if path.isfile(gbk_name):
                gbk_name = str(gbk_name).split('(')[0] + f'({x}).' + str(gbk_name).split(')')[1]
                x += 1
            else:
                break
    rename(name, gbk_name)
    process_ppt(gbk_name)


def main(choice1):  # 主函数
    choice_link = url_all + get_choice_link(choice1)
    web_for_page = req_for_web(choice_link)
    for a in range(1, int(get_the_max_page(web_for_page))):
        add_page = choice_link + format_link(web_for_page, a)
        file_list = file_download(add_page)
        with tqdm(total=len(file_list),
                  bar_format='第%d页下载中:{percentage:3.0f}%%|{bar}|{n}/{total}[{rate_fmt}{postfix}]' % a) as bar:
            # 进度条实现
            for link in file_list:
                req_file = get(link[0], headers=headers).content
                data = io.BytesIO()
                data.write(req_file)
                zip_file = zipfile.ZipFile(data)
                for names in zip_file.namelist():  # 文件解压
                    if '.ppt' in names:
                        zip_file.extract(names)
                        auto_decode(names)
                zip_file.close()
                bar.update(1)


if __name__ == "__main__":
    main(choice)
